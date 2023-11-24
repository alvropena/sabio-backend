from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.responses import JSONResponse
from pptx import Presentation
from io import BytesIO
from fastapi.middleware.cors import CORSMiddleware

app = FastAPI()

# Allow all origins in development; adjust the origins parameter accordingly in production
origins = ["http://localhost:3000"]  # Replace with your Next.js frontend URL

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.post("/api/extract/")
async def extract_text_from_pptx(file: UploadFile = File(...)):
    try:
        pptx_file = await file.read()
        text_data = extract_text_from_pptx(pptx_file)
        return JSONResponse(content=text_data, status_code=201)
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))


def extract_text_from_pptx(pptx_file):
    presentation = Presentation(BytesIO(pptx_file))
    text_data = []

    for i, slide in enumerate(presentation.slides):
        text = '\n'.join(
            [shape.text for shape in slide.shapes if hasattr(shape, 'text')])
        text_data.append({'slide_number': i + 1, 'text': text})

    return text_data
